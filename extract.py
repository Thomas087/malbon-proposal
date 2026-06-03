#!/usr/bin/env python3
import json, os, math
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = "0601_Malbon_Proposal_Structure_v3 (hidden removed).pptx"
IMG_ROOT = "images"
EMU_PER_IN = 914400.0

prs = Presentation(PPTX)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

os.makedirs(IMG_ROOT, exist_ok=True)

def emu_to_in(v):
    if v is None: return None
    return round(v / EMU_PER_IN, 3)

def gcd_ratio(w, h):
    if not w or not h: return None
    g = math.gcd(int(w), int(h))
    return f"{int(w/g)}:{int(h/g)}"

def pos(shape):
    try:
        return {
            "left_in": emu_to_in(shape.left),
            "top_in": emu_to_in(shape.top),
            "width_in": emu_to_in(shape.width),
            "height_in": emu_to_in(shape.height),
        }
    except Exception:
        return {}

def extract_text(tf):
    paras = []
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs)
        if not text and p.text:
            text = p.text
        if text.strip() == "":
            continue
        # gather formatting hints
        bold = any(r.font.bold for r in p.runs if r.font.bold is not None)
        sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        paras.append({
            "level": p.level,
            "text": text,
            "bold": bold,
            "font_pt": max(sizes) if sizes else None,
        })
    return paras

def handle_picture(shape, slide_dir, slide_no, counter):
    img = shape.image
    ext = img.ext
    px_w, px_h = img.size  # native pixels
    fname = f"slide{slide_no:02d}_img{counter}.{ext}"
    out_path = os.path.join(slide_dir, fname)
    with open(out_path, "wb") as f:
        f.write(img.blob)
    return {
        "type": "picture",
        "name": shape.name,
        "file": out_path,
        "native_px": [px_w, px_h],
        "native_aspect_ratio": gcd_ratio(px_w, px_h),
        "native_aspect_decimal": round(px_w/px_h, 3) if px_h else None,
        "display": pos(shape),
        "display_aspect_decimal": round(shape.width/shape.height, 3) if shape.height else None,
    }

def walk(shape, slide_dir, slide_no, counter):
    """Returns (node, new_counter)"""
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        children = []
        for sub in shape.shapes:
            node, counter = walk(sub, slide_dir, slide_no, counter)
            if node:
                children.append(node)
        return ({
            "type": "group",
            "name": shape.name,
            "position": pos(shape),
            "children": children,
        } if children else None), counter
    if st == MSO_SHAPE_TYPE.PICTURE or (hasattr(shape, "image") and st == 13):
        try:
            counter += 1
            return handle_picture(shape, slide_dir, slide_no, counter), counter
        except Exception as e:
            return {"type":"picture_error","name":shape.name,"error":str(e)}, counter
    if shape.has_table:
        tbl = shape.table
        rows = []
        for r in tbl.rows:
            rows.append([c.text for c in r.cells])
        return {"type":"table","name":shape.name,"position":pos(shape),"rows":rows}, counter
    if shape.has_text_frame:
        paras = extract_text(shape.text_frame)
        if paras:
            return {"type":"text","name":shape.name,"position":pos(shape),"paragraphs":paras}, counter
        return None, counter
    # try picture fallback
    try:
        if hasattr(shape, "image"):
            counter += 1
            return handle_picture(shape, slide_dir, slide_no, counter), counter
    except Exception:
        pass
    # other shapes (could be placeholder/autoshape with no text)
    return None, counter

slides_out = []
for i, slide in enumerate(prs.slides, start=1):
    slide_dir = os.path.join(IMG_ROOT, f"slide{i:02d}")
    os.makedirs(slide_dir, exist_ok=True)
    counter = 0
    shapes_nodes = []
    for shape in slide.shapes:
        node, counter = walk(shape, slide_dir, i, counter)
        if node:
            shapes_nodes.append(node)
    # remove empty image dir
    if not os.listdir(slide_dir):
        os.rmdir(slide_dir)
    slides_out.append({
        "slide_number": i,
        "shapes": shapes_nodes,
    })

out = {
    "source_file": PPTX,
    "slide_size_in": [emu_to_in(SLIDE_W), emu_to_in(SLIDE_H)],
    "slide_count": len(slides_out),
    "slides": slides_out,
}
with open("malbon_raw_extract.json", "w") as f:
    json.dump(out, f, indent=2, ensure_ascii=False)

print("slides:", len(slides_out))
print("slide size in:", out["slide_size_in"])
# quick stats
total_imgs = sum(1 for s in slides_out for _ in [1])
print("written malbon_raw_extract.json")
